# -*- coding: utf-8 -*-
"""
HTML 프레젠테이션을 PowerPoint로 변환하는 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_portfolio_presentation():
    """React Portfolio 프레젠테이션 생성"""
    
    # 프레젠테이션 객체 생성
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 색상 정의 (HTML에서 사용된 색상)
    BG_DEEP = RGBColor(2, 6, 23)
    BG_OCEAN = RGBColor(15, 23, 42)
    BG_ROYAL = RGBColor(23, 37, 84)
    TEXT_WHITE = RGBColor(248, 250, 252)
    TEXT_BLUE_GRAY = RGBColor(148, 163, 184)
    TEXT_ACCENT = RGBColor(56, 189, 248)
    
    # 슬라이드 0: INTRO
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    
    # REACT 타이틀
    title_texts = [
        ("R", "eact"),
        ("E", "xperience"),
        ("A", "rchive"),
        ("C", "reative"),
        ("T", "imeline")
    ]
    
    top = Inches(1.5)
    for head, tail in title_texts:
        # Head letter
        txBox = slide.shapes.add_textbox(Inches(3), top, Inches(10), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run_head = p.add_run()
        run_head.text = head
        run_head.font.size = Pt(72)
        run_head.font.bold = True
        run_head.font.color.rgb = TEXT_WHITE
        
        # Tail text
        run_tail = p.add_run()
        run_tail.text = tail
        run_tail.font.size = Pt(42)
        run_tail.font.color.rgb = TEXT_BLUE_GRAY
        
        top += Inches(0.9)
    
    # 서브타이틀
    txBox = slide.shapes.add_textbox(Inches(3), top + Inches(0.5), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    tf.text = "React Portfolio by Sunhayoung\n2025.12"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = TEXT_ACCENT
    
    # 슬라이드 1: Table of Contents
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Table of Contents"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "목차"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 목차 항목
    toc_items = [
        ("01", "Overview", "프로젝트 개요"),
        ("02", "Features", "주요 기능"),
        ("03", "Tech Stack · Structure", "기술 스택과 프로젝트 구조"),
        ("04", "Solving", "문제 해결 과정"),
        ("05", "Timeline", "개발 로그"),
        ("06", "Outcomes", "성과 및 회고")
    ]
    
    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(4)
    height = Inches(1.8)
    
    for i, (num, title, desc) in enumerate(toc_items):
        col = i % 3
        row = i // 3
        
        box_left = left + col * (width + Inches(0.5))
        box_top = top + row * (height + Inches(0.3))
        
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            box_left, box_top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 58, 138)
        shape.line.color.rgb = RGBColor(96, 165, 250)
        
        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        
        # Number
        p = text_frame.paragraphs[0]
        p.text = num
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Title
        p = text_frame.add_paragraph()
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Description
        p = text_frame.add_paragraph()
        p.text = desc
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 슬라이드 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Overview"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "리액트 기반 인터랙티브 포트폴리오 웹사이트 구현"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 3개 카드
    cards = [
        ("🎯 Goal", "단순 클론 코딩을 넘어,\n데이터 흐름과 상태 관리를\n직접 설계하고 구현하는 것."),
        ("📅 Duration", "2025.12.15 (월) ~ 12.24 (수)\n9 Days Sprint | 지속 개발 중"),
        ("⚡ Key Result", "Functional Components & Hooks\n실무적 활용 및 비즈니스 로직 추상화\nRedux Toolkit으로 복잡한 상태 관리 체계화")
    ]
    
    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(4)
    height = Inches(4)
    
    for i, (title, content) in enumerate(cards):
        box_left = left + i * (width + Inches(0.5))
        
        shape = slide.shapes.add_shape(1, box_left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 58, 138)
        shape.line.color.rgb = RGBColor(96, 165, 250)
        
        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_right = Inches(0.3)
        
        # Title
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Content
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.space_before = Pt(10)
    
    # 슬라이드 3: Core Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Features"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "사용자 경험(UX)을 최우선으로 고려한 기능 구현"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 6개 features
    features = [
        ("✅ Home Page", "디자인물 카드 그리드 레이아웃\n썸네일 + 제목 + 태그 구성"),
        ("✅ Category Filter", "Redux 기반 실시간 필터링\n즉시 반응하는 UI 업데이트"),
        ("✅ Search", "제목/태그 실시간 검색 구현\nuseCallback 최적화"),
        ("✅ Sort Logic", "최신/오래된순/가나다 정렬\nRedux Selector 활용"),
        ("✅ Detail Zoom", "95vh 고해상도 확대 뷰어\nHyper-Sensitive Zoom"),
        ("✅ Theme & Responsive", "다크/라이트 모드 (Context API)\n완벽한 모바일-데스크톱 대응")
    ]
    
    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(4)
    height = Inches(2.5)
    
    for i, (title, content) in enumerate(features):
        col = i % 3
        row = i // 3
        
        box_left = left + col * (width + Inches(0.5))
        box_top = top + row * (height + Inches(0.3))
        
        shape = slide.shapes.add_shape(1, box_left, box_top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 58, 138)
        shape.line.color.rgb = RGBColor(96, 165, 250)
        
        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_left = Inches(0.3)
        
        # Title
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(21)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Content
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(203, 213, 225)
    
    # 슬라이드 4: Tech Stack & Structure
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Structure & Tech Stack"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "견고한 아키텍처와 기술 선정"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_BLUE_GRAY
    
    # Left side - Tech stack (3 cards)
    tech_cards = [
        ("Frontend & Build", "- React 18: 최신 Hooks API\n- Vite: 빠른 HMR & 번들링\n- Router v6: SPA 동적 라우팅"),
        ("State Management", "- Redux Toolkit: Slice 패턴 관리\n  (portfolio, theme, ui slices)\n- Context API: Theme 전역 관리"),
        ("Advanced Patterns", "- Custom Hooks: 로직 추상화\n- Lazy Loading: 초기 로딩 최적화")
    ]
    
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(7)
    height = Inches(1.8)
    
    for i, (title, content) in enumerate(tech_cards):
        box_top = top + i * (height + Inches(0.2))
        
        shape = slide.shapes.add_shape(1, left, box_top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 58, 138)
        shape.line.color.rgb = RGBColor(96, 165, 250)
        
        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_left = Inches(0.3)
        
        # Title
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Content
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(203, 213, 225)
    
    # Right side - Project structure
    structure_text = """📂 Project Tree

src
 ├── components
 │   ├── layout (Header, Footer)
 │   ├── portfolio (Card, List)
 │   └── ui (Buttons, Scroll)
 ├── pages
 │   ├── Home, About, Work
 │   └── Detail, Contact
 ├── store
 │   └── slices (portfolio, theme)
 ├── theme (Context, Hooks)
 ├── utils (Images, Helpers)
 └── App.jsx"""
    
    shape = slide.shapes.add_shape(1, Inches(8.5), Inches(2.5), Inches(6.5), Inches(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
    shape.line.color.rgb = RGBColor(59, 130, 246)
    
    text_frame = shape.text_frame
    text_frame.margin_top = Inches(0.3)
    text_frame.margin_left = Inches(0.3)
    
    p = text_frame.paragraphs[0]
    p.text = structure_text
    p.font.size = Pt(13)
    p.font.name = "Consolas"
    p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 슬라이드 5: Problem & Solving
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Problem & Solving"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CODE LEVEL Trouble Shooting"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 4개의 문제 해결
    problems = [
        ("🔴 1. Git Pages Asset 404",
         "Base URL 경로 불일치로 리소스 로드 실패",
         "Before: ${baseUrl}${cleanPath}",
         "After: path.startsWith('/') ? path.slice(1) : path",
         "순수 상대 경로로 변환하여 해결"),
        ("🔴 2. CSS Import Order",
         "컴포넌트 로드 후 스타일 적용 실패",
         "const Comp = () => {...}\nimport './Comp.css'; // ❌",
         "import './Comp.css'; // ✅\nconst Comp = () => {...}",
         "Import Hoisting으로 우선순위 확보"),
        ("🔴 3. Redux Infinite Render",
         "filter()가 매번 새 배열을 반환",
         "useSelector(state => state.items.filter(...))",
         "createSelector(items, filters, (i, f) => ... )",
         "Reselect 메모이제이션으로 참조 유지"),
        ("🔴 4. Vite HMR Frozen",
         "Windows 환경 파일 감시 충돌",
         "// No Config (Default)",
         "server: { hmr: { overlay: false }, watch: {...} }",
         "Config 명시적 지정으로 안정화")
    ]
    
    left = Inches(1)
    top = Inches(2.3)
    width = Inches(7)
    height = Inches(3.0)
    
    for i, (title, desc, before, after, solution) in enumerate(problems):
        col = i % 2
        row = i // 2
        
        box_left = left + col * (width + Inches(0.5))
        box_top = top + row * (height + Inches(0.2))
        
        shape = slide.shapes.add_shape(1, box_left, box_top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
        shape.line.color.rgb = RGBColor(59, 130, 246)
        
        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.word_wrap = True
        
        # Title
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(248, 113, 113)
        
        # Description
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(203, 213, 225)
        
        # Before
        p = text_frame.add_paragraph()
        p.text = before[:50] + "..." if len(before) > 50 else before
        p.font.size = Pt(9)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(248, 113, 113)
        
        # After (Solution)
        p = text_frame.add_paragraph()
        p.text = "🟢 Solution"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(74, 222, 128)
        
        p = text_frame.add_paragraph()
        p.text = after[:50] + "..." if len(after) > 50 else after
        p.font.size = Pt(9)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(74, 222, 128)
        
        # Solution description
        p = text_frame.add_paragraph()
        p.text = solution
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 슬라이드 6: Development Log
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Development Log"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "개발 과정 7단계"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 7 phases
    phases = [
        ("Phase 1", "Init & Basics", "Vite 템플릿 세팅\nRedux Store 구성\nRouter 설계\n폴더 구조 정의"),
        ("Phase 2", "Layout & Theme", "Main Layout 구현\nCSS Variables\nTheme Context\nDark/Light Mode"),
        ("Phase 3", "Core Components", "Portfolio Card/List\nHero Section\nFilter & Search UI\nGrid Layout"),
        ("Phase 4", "Detail & Zoom", "Dynamic Routing\n2-Stage Zoom\nImage Gallery\nuseParams"),
        ("Phase 5", "UI/UX Polish", "Hybrid Typography\nGlassmorphism\nHover Effects\nMicro Interactions"),
        ("Phase 6", "Responsive", "Media Queries\nAccessiblity\nLazy Loading"),
        ("Phase 7", "Deploy", "Code Formatting\nREADME 작성\nGitHub Pages")
    ]
    
    left = Inches(0.8)
    top = Inches(2.3)
    width = Inches(3.4)
    height = Inches(2.8)
    
    for i, (phase, title, content) in enumerate(phases):
        col = i % 4
        row = i // 4
        
        box_left = left + col * (width + Inches(0.2))
        box_top = top + row * (height + Inches(0.3))
        
        shape = slide.shapes.add_shape(1, box_left, box_top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 58, 138)
        shape.line.color.rgb = RGBColor(96, 165, 250)
        shape.line.width = Pt(2)
        
        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_left = Inches(0.2)
        
        # Phase number
        p = text_frame.paragraphs[0]
        p.text = phase
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_ACCENT
        
        # Title
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Content
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(203, 213, 225)
    
    # 슬라이드 7: Outcomes & Growth
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Outcomes & Growth"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "프로젝트를 통해 얻은 것"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 4 outcomes
    outcomes = [
        ("⚛️ React Ecosystem",
         "- Hooks: 의존성 배열 원리 및 클로저 문제 해결\n- Custom Hook: 로직 재사용성 및 구조화\n- Performance: memo, useMemo 적절한 사용"),
        ("🏗 State Architecture",
         "- Redux vs Local: 전역/지역 상태의 명확한 구분\n- Selector Pattern: 상태 조회 최적화\n- Slice Design: 도메인별 상태 분리"),
        ("🎨 Styling Philosophy",
         "- CSS Variables: 동적 테마 유연성 확보\n- Modern Layout: Grid/Flexbox 반응형 설계\n- Zero Runtime: 순수 CSS로 오버헤드 최소화"),
        ("🛣 SPA Strategy",
         "- URL Sync: URL ↔ UI 상태 동기화\n- History: 브라우저 네비게이션 완벽 지원\n- Deep Linking: 직접 접근 가능한 경로 설계")
    ]
    
    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(6.5)
    height = Inches(2.7)
    
    for i, (title, content) in enumerate(outcomes):
        col = i % 2
        row = i // 2
        
        box_left = left + col * (width + Inches(0.5))
        box_top = top + row * (height + Inches(0.3))
        
        shape = slide.shapes.add_shape(1, box_left, box_top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 58, 138)
        shape.line.color.rgb = RGBColor(96, 165, 250)
        
        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_left = Inches(0.3)
        
        # Title
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Content
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(203, 213, 225)
    
    # 슬라이드 8: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    
    # Thank You
    title_box = slide.shapes.add_textbox(Inches(3), Inches(2), Inches(10), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    # 메시지
    msg_box = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(10), Inches(2))
    tf = msg_box.text_frame
    p = tf.paragraphs[0]
    p.text = "단순한 구현을 넘어, 가치 있는 사용자 경험을 설계합니다.\n끊임없이 고민하고 성장하는 개발자가 되겠습니다."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_BLUE_GRAY
    
    # 링크
    link_box = slide.shapes.add_textbox(Inches(5), Inches(6.5), Inches(6), Inches(1))
    tf = link_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🌐 Live Demo: hayohio-bit.github.io/hayohio-portfolio/\n🔗 GitHub: github.com/hayohio-bit/hayohio-portfolio"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_ACCENT
    
    # 저장
    prs.save('portfolio_presentation.pptx')
    print("✅ PowerPoint 파일이 성공적으로 생성되었습니다: portfolio_presentation.pptx")

if __name__ == "__main__":
    create_portfolio_presentation()
